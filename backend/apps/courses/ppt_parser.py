"""Teaching-file parsing and catalog extraction.

The API kept the early ``ppts`` naming, but course materials can now be PPT,
PDF, Word, text, Markdown, or CSV.  This module extracts text blocks first, then
uses an AI provider to identify a clean catalog tree.  If AI is unavailable or
returns invalid JSON, a conservative rule-based fallback is used.
"""
from __future__ import annotations

import csv
import json
import os
import re
import shutil
from pathlib import Path


SUPPORTED_EXTENSIONS = {".ppt", ".pptx", ".pdf", ".doc", ".docx", ".txt", ".md", ".csv"}

META_PATTERNS = (
    "学年",
    "学期",
    "课程设计",
    "授课计划",
    "课程名称",
    "任课教师",
    "教师姓名",
    "专业",
    "班级",
    "学时",
    "学分",
    "考核",
    "教材",
    "日期",
)


def parse_ppt_pages(file_path: str) -> list[dict]:
    """Backward-compatible alias used by existing views."""
    return parse_teaching_file_pages(file_path)


def render_presentation_slide_images(file_path: str, resource_id: int | str | None = None) -> dict[int, str]:
    """Render PPT/PPTX slides to PNG images and return {page: media_url}.

    On Windows this uses the installed Microsoft PowerPoint COM automation.
    If rendering is unavailable, callers still keep text extraction working.
    """
    if not file_path or not os.path.exists(file_path):
        return {}
    ext = os.path.splitext(file_path)[1].lower()
    if ext not in (".ppt", ".pptx"):
        return {}

    try:
        from django.conf import settings
        import pythoncom
        import win32com.client
    except Exception:
        return {}

    rel_dir = Path("ppt_pages") / str(resource_id or Path(file_path).stem)
    out_dir = Path(settings.MEDIA_ROOT) / rel_dir
    if out_dir.exists():
        shutil.rmtree(out_dir, ignore_errors=True)
    out_dir.mkdir(parents=True, exist_ok=True)

    app = None
    presentation = None
    rendered: dict[int, str] = {}
    pythoncom.CoInitialize()
    try:
        app = win32com.client.Dispatch("PowerPoint.Application")
        try:
            app.DisplayAlerts = 0
        except Exception:
            pass
        presentation = app.Presentations.Open(os.path.abspath(file_path), True, False, False)
        width = int(getattr(presentation.PageSetup, "SlideWidth", 1280) * 2)
        height = int(getattr(presentation.PageSetup, "SlideHeight", 720) * 2)
        for index, slide in enumerate(presentation.Slides, start=1):
            image_name = f"slide_{index:03d}.png"
            image_path = out_dir / image_name
            slide.Export(str(image_path), "PNG", width, height)
            if image_path.exists():
                rendered[index] = _media_url(str(rel_dir / image_name))
    except Exception:
        return rendered
    finally:
        if presentation is not None:
            try:
                presentation.Close()
            except Exception:
                pass
        if app is not None:
            try:
                app.Quit()
            except Exception:
                pass
        pythoncom.CoUninitialize()
    return rendered


def attach_slide_images(pages: list[dict], images: dict[int, str]) -> list[dict]:
    if not images:
        return pages
    for page in pages:
        page_no = page.get("page")
        if page_no in images:
            page["image"] = images[page_no]
    return pages


def _media_url(relative_path: str) -> str:
    from django.conf import settings

    base = settings.MEDIA_URL if settings.MEDIA_URL.startswith("/") else f"/{settings.MEDIA_URL}"
    return f"{base.rstrip('/')}/{relative_path.replace(os.sep, '/')}"


def parse_teaching_file_pages(file_path: str) -> list[dict]:
    if not file_path or not os.path.exists(file_path):
        return []

    ext = os.path.splitext(file_path)[1].lower()
    if ext not in SUPPORTED_EXTENSIONS:
        return []
    if ext in (".ppt", ".pptx"):
        return _read_presentation(file_path)
    if ext == ".pdf":
        return _read_pdf(file_path)
    if ext in (".doc", ".docx"):
        return _read_docx(file_path)
    if ext in (".txt", ".md"):
        return _read_text_outline(file_path)
    if ext == ".csv":
        return _read_csv(file_path)
    return []


def pages_to_catalog_tree(pages: list[dict]) -> list[dict]:
    """Return catalog tree using AI first, then a rule fallback."""
    if not pages:
        return []

    ai_tree = _ai_extract_catalog_tree(pages)
    if ai_tree:
        return ai_tree
    return _fallback_catalog_tree(pages)


def pages_to_chapter_tree(pages: list[dict]) -> list[dict]:
    """Return only top-level chapter titles for initial course catalogs."""
    if not pages:
        return []
    # 授课计划/教学大纲通常是结构清晰的“第X章”列表：优先用规则逐字抽取——
    # 快速、稳定、忠于原文，也避免大模型不可用时长时间阻塞请求（前端 60s 超时）。
    rule_chapters = _fallback_chapter_tree(pages)
    if rule_chapters:
        return rule_chapters
    # 规则没抽到（如无“第X章”标记的零散文本）时，再求助 AI 归纳。
    tree = pages_to_catalog_tree(pages)
    return _chapter_tree_from_nodes(tree)


def _ai_extract_catalog_tree(pages: list[dict]) -> list[dict]:
    text = _pages_to_prompt_text(pages)
    if not text.strip():
        return []

    try:
        from apps.ai.providers.factory import get_provider

        provider = get_provider()
        # 纯 Mock（未配置真实大模型）无法真正读取文件内容，返回的是与上传文件无关的
        # 示例目录。此时直接放弃 AI 抽取，改由下游规则解析真实文本。
        if getattr(provider, "name", "") == "mock":
            return []
        messages = [
            {
                "role": "system",
                "content": (
                    "你是课程目录抽取器，只做【逐字抽取】，严禁改写、合并、总结、虚构。"
                    "只输出 JSON 数组，不要任何解释。"
                ),
            },
            {
                "role": "user",
                "content": (
                    "从下面授课计划文本中提取课程目录，严格遵守：\n"
                    "1) 章：以“第X章”开头的行（如“第1章 Java语言概述”）为章标题，"
                    "必须逐字保留原文，原文有几章就提取几章，不得少、不得多、不得改字（不要把“第1章”写成“第一章”）。\n"
                    "2) 节：形如“X.Y ...”的行（如“3.3 类的封装性”）为节，归入其所属章的 children。\n"
                    "3) 忽略：实验/习题/教材/学时/课时/黑板/多媒体/计算机/实验报告 等行，以及页码、课程名称/专业/教师/班级/学期等封面元信息。\n"
                    "4) 绝对不要自行发明原文没有的章节标题，只能来自下面文本。\n"
                    "输出示例：[{\"title\":\"第1章 xxx\",\"children\":[{\"title\":\"1.1 xxx\"}]}]\n\n"
                    f"授课计划文本：\n{text[:12000]}"
                ),
            },
        ]
        # fallback_to_mock=False：真实模型不可用时直接抛错，避免拿到与上传文件无关的
        # 「第一章 绪论/第二章 基础知识」占位目录；异常由下方 except 捕获后走规则解析。
        # retries/timeout 收紧：网络挂起时也能快速失败，不拖垮前端 60s 请求预算。
        raw = provider.chat(messages, temperature=0.1, fallback_to_mock=False, retries=1, timeout=15)
    except Exception:
        return []

    data = _safe_json(raw)
    if not isinstance(data, list):
        return []
    return _normalize_tree(data)


def _chapter_tree_from_nodes(nodes: list[dict]) -> list[dict]:
    chapters: list[dict] = []
    seen = set()

    def collect(items: list[dict]) -> None:
        for item in items or []:
            title = _extract_chapter_title(str(item.get("title", "")))
            key = re.sub(r"\s+", "", title)
            if title and key not in seen:
                chapters.append({"title": title, "children": []})
                seen.add(key)
            collect(item.get("children") or [])

    collect(nodes)
    return chapters


def _fallback_chapter_tree(pages: list[dict]) -> list[dict]:
    chapters: list[dict] = []
    seen = set()
    for page in pages:
        for text in (page.get("title", ""), page.get("body", "")):
            for line in _split_candidate_lines(text):
                title = _extract_chapter_title(line)
                key = re.sub(r"\s+", "", title)
                if title and key not in seen:
                    chapters.append({"title": title, "children": []})
                    seen.add(key)
    return chapters


def _fallback_catalog_tree(pages: list[dict]) -> list[dict]:
    candidates: list[str] = []
    for page in pages:
        for text in (page.get("title", ""), page.get("body", "")):
            candidates.extend(_split_candidate_lines(text))

    cleaned: list[str] = []
    seen = set()
    for line in candidates:
        line = _clean_title(line)
        if not line or line in seen or _is_meta_line(line):
            continue
        if _looks_like_catalog_line(line):
            cleaned.append(line)
            seen.add(line)

    if not cleaned:
        for page in pages:
            title = _clean_title(page.get("title", ""))
            if title and title not in seen and not _is_meta_line(title):
                cleaned.append(title)
                seen.add(title)

    tree: list[dict] = []
    current_chapter: dict | None = None
    for title in cleaned[:80]:
        node = {"title": title, "children": []}
        if _looks_like_section(title) and current_chapter:
            current_chapter["children"].append(node)
        else:
            tree.append(node)
            current_chapter = node
    return tree


def _read_presentation(file_path: str) -> list[dict]:
    try:
        from pptx import Presentation

        prs = Presentation(file_path)
    except Exception:
        return []

    pages = []
    for idx, slide in enumerate(prs.slides, start=1):
        title = ""
        bodies = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text.strip()
            if not text:
                continue
            is_title = False
            try:
                is_title = bool(shape.is_placeholder and shape.placeholder_format.idx == 0)
            except Exception:
                pass
            if is_title and not title:
                title = text
            else:
                bodies.append(text)
        if not title and bodies:
            title = bodies[0][:80]
        pages.append({"page": idx, "title": title, "body": "\n".join(bodies), "source": "slide"})
    return pages


def _read_pdf(file_path: str) -> list[dict]:
    try:
        from pypdf import PdfReader

        reader = PdfReader(file_path)
        pages = []
        for idx, page in enumerate(reader.pages, start=1):
            text = (page.extract_text() or "").strip()
            if not text:
                continue
            pages.append({"page": idx, "title": _first_line(text)[:80], "body": text, "source": "pdf"})
        return pages
    except Exception:
        return []


def _read_docx(file_path: str) -> list[dict]:
    try:
        from docx import Document

        doc = Document(file_path)
    except Exception:
        return []

    pages: list[dict] = []

    for table_index, table in enumerate(doc.tables, start=1):
        rows = [[_cell_text(cell) for cell in row.cells] for row in table.rows]
        if not rows:
            continue
        headers = rows[0]
        content_indexes = _content_column_indexes(headers)
        if not content_indexes:
            content_indexes = list(range(len(headers)))
        for row_index, row in enumerate(rows[1:], start=1):
            values = [row[i] for i in content_indexes if i < len(row) and row[i]]
            if not values:
                continue
            title = _pick_best_title(values)
            body = "\n".join(v for v in row if v)
            if title and not _is_meta_line(title):
                pages.append(
                    {
                        "page": len(pages) + 1,
                        "title": title,
                        "body": body,
                        "source": f"table-{table_index}-{row_index}",
                    }
                )

    paragraph_pages = _read_docx_paragraphs(doc)
    if pages:
        pages.extend(p for p in paragraph_pages if _looks_like_catalog_line(p.get("title", "")))
        return pages
    return paragraph_pages


def _read_docx_paragraphs(doc) -> list[dict]:
    pages: list[dict] = []
    current: dict | None = None
    fallback_lines: list[str] = []

    for paragraph in doc.paragraphs:
        text = paragraph.text.strip()
        if not text:
            continue
        style_name = (getattr(paragraph.style, "name", "") or "").lower()
        is_heading = style_name.startswith("heading") or _looks_like_catalog_line(text)
        if is_heading and not _is_meta_line(text):
            if current:
                pages.append(current)
            current = {"page": len(pages) + 1, "title": text[:100], "body": "", "source": "paragraph"}
        elif current:
            current["body"] = (current["body"] + "\n" + text).strip()
        else:
            fallback_lines.append(text)
    if current:
        pages.append(current)

    if pages:
        return pages
    return _lines_to_pages(fallback_lines)


def _read_text_outline(file_path: str) -> list[dict]:
    text = _read_text(file_path)
    if not text:
        return []
    return _lines_to_pages(text.splitlines())


def _read_csv(file_path: str) -> list[dict]:
    try:
        with open(file_path, newline="", encoding="utf-8-sig") as f:
            rows = list(csv.reader(f))
    except UnicodeDecodeError:
        with open(file_path, newline="", encoding="gbk") as f:
            rows = list(csv.reader(f))
    except Exception:
        return []

    pages = []
    headers = rows[0] if rows else []
    content_indexes = _content_column_indexes(headers) or list(range(len(headers)))
    for idx, row in enumerate(rows[1:] if headers else rows, start=1):
        values = [row[i].strip() for i in content_indexes if i < len(row) and row[i].strip()]
        if not values:
            continue
        title = _pick_best_title(values)
        if title and not _is_meta_line(title):
            pages.append({"page": idx, "title": title, "body": " / ".join(row), "source": "csv"})
    return pages


def _content_column_indexes(headers: list[str]) -> list[int]:
    keys = ("教学内容", "授课内容", "课程内容", "章节", "章", "节", "项目", "任务", "知识点", "主要内容")
    indexes = []
    for index, header in enumerate(headers):
        compact = re.sub(r"\s+", "", header or "")
        if any(key in compact for key in keys):
            indexes.append(index)
    return indexes


def _pick_best_title(values: list[str]) -> str:
    lines = []
    for value in values:
        lines.extend(_split_candidate_lines(value))
    catalog_like = [line for line in lines if _looks_like_catalog_line(line) and not _is_meta_line(line)]
    if catalog_like:
        return catalog_like[0][:100]
    meaningful = [line for line in lines if len(line) >= 4 and not _is_meta_line(line)]
    return meaningful[0][:100] if meaningful else ""


def _pages_to_prompt_text(pages: list[dict]) -> str:
    parts = []
    for item in pages[:160]:
        title = item.get("title", "")
        body = item.get("body", "")
        source = item.get("source", "")
        parts.append(f"[{source}] {title}\n{body}".strip())
    return "\n\n---\n\n".join(parts)


def _safe_json(text: str):
    if not text:
        return None
    cleaned = text.strip()
    cleaned = re.sub(r"^```(?:json)?", "", cleaned).strip()
    cleaned = re.sub(r"```$", "", cleaned).strip()
    try:
        return json.loads(cleaned)
    except Exception:
        pass
    match = re.search(r"(\[[\s\S]*\]|\{[\s\S]*\})", cleaned)
    if not match:
        return None
    try:
        return json.loads(match.group(1))
    except Exception:
        return None


def _normalize_tree(data: list) -> list[dict]:
    tree = []
    for item in data:
        if not isinstance(item, dict):
            continue
        title = _clean_title(str(item.get("title", "")))
        if not title or _is_meta_line(title):
            continue
        children = []
        for child in item.get("children") or []:
            if not isinstance(child, dict):
                continue
            child_title = _clean_title(str(child.get("title", "")))
            if child_title and not _is_meta_line(child_title):
                children.append({"title": child_title, "children": []})
        tree.append({"title": title, "children": children})
    return tree


def _cell_text(cell) -> str:
    return "\n".join(p.text.strip() for p in cell.paragraphs if p.text.strip())


def _read_text(file_path: str) -> str:
    for encoding in ("utf-8", "utf-8-sig", "gbk", "latin-1"):
        try:
            return Path(file_path).read_text(encoding=encoding)
        except UnicodeDecodeError:
            continue
        except Exception:
            return ""
    return ""


def _lines_to_pages(lines: list[str], chunk_size: int = 12) -> list[dict]:
    cleaned = [_clean_title(line) for line in lines if _clean_title(line)]
    pages = []
    current = None
    for line in cleaned:
        if _is_meta_line(line):
            continue
        if _looks_like_catalog_line(line):
            if current:
                pages.append(current)
            current = {"page": len(pages) + 1, "title": line[:100], "body": "", "source": "text"}
        elif current:
            current["body"] = (current["body"] + "\n" + line).strip()
    if current:
        pages.append(current)
    if pages:
        return pages

    for start in range(0, len(cleaned), chunk_size):
        chunk = [line for line in cleaned[start : start + chunk_size] if not _is_meta_line(line)]
        if not chunk:
            continue
        pages.append({"page": len(pages) + 1, "title": chunk[0][:100], "body": "\n".join(chunk), "source": "text"})
    return pages


def _split_candidate_lines(text: str) -> list[str]:
    if not text:
        return []
    # 垂直制表符/换页符是标题内的“软换行”（如 PPT 标题“第一章\x0bJava概述”），
    # 视作空格而非分行，避免一个标题被 splitlines 拆成两段。
    text = text.replace("\v", " ").replace("\f", " ")
    normalized = re.sub(r"[；;]\s*", "\n", text)
    normalized = re.sub(r"\s{2,}", "\n", normalized)
    return [_clean_title(line) for line in normalized.splitlines() if _clean_title(line)]


def _clean_title(text: str) -> str:
    text = re.sub(r"\s+", " ", text or "").strip()
    text = re.sub(r"^[\-\*●•·]\s*", "", text)
    return text.strip(" ：:，,。")


def _extract_chapter_title(text: str) -> str:
    text = _clean_title(text)
    if not text or _is_meta_line(text):
        return ""

    # 仅识别“第X章 / 第X讲 / 第X单元”这类章目录（与前端一致、也符合本功能既定范围）。
    # 不再把“1. xxx / 一、xxx”这类课件正文里的编号条目误当成章，避免 PPT 正文被大量误识别。
    patterns = (
        r"(第\s*[一二三四五六七八九十百千万\d]+\s*(?:章|讲|单元)\s*[^；;\n]*)",
    )
    for pattern in patterns:
        match = re.search(pattern, text, flags=re.IGNORECASE)
        if not match:
            continue
        title = _clean_title(match.group(1))
        title = re.split(
            r"\s+(?:\d+\.\d+|第\s*[一二三四五六七八九十百千万\d]+\s*节|实验\d*|习题|作业|学时|课时|考核|教材|参考|备注)\b",
            title,
            maxsplit=1,
        )[0]
        title = re.sub(r"^(第)\s+([一二三四五六七八九十百千万\d]+)\s+(章|讲|单元)", r"\1\2\3", title)
        title = re.sub(r"^(第[一二三四五六七八九十百千万\d]+(?:章|讲|单元))(?=\S)", r"\1 ", title)
        title = _clean_title(title)
        if title and not _is_meta_line(title) and len(title) <= 60:
            return title
    return ""


def _first_line(text: str) -> str:
    return next((line.strip() for line in text.splitlines() if line.strip()), "")


def _is_meta_line(text: str) -> bool:
    compact = re.sub(r"\s+", "", text or "")
    if not compact:
        return True
    if len(compact) > 120:
        return True
    if re.search(r"\d{4}\s*/\s*\d{4}\s*学年", compact):
        return True
    return any(pattern in compact for pattern in META_PATTERNS) and not _looks_like_catalog_line(compact)


def _looks_like_catalog_line(text: str) -> bool:
    text = _clean_title(text)
    if not text or len(text) > 100:
        return False
    patterns = (
        r"^第[一二三四五六七八九十百\d]+[章节讲课单元]",
        r"^\d+([.、]\d+)*[.、\s]",
        r"^[一二三四五六七八九十]+[、.]\s*",
        r".*(绪论|概述|基础|原理|实践|案例|项目|实验|复习|总结|设计|开发|应用).*",
    )
    return any(re.search(pattern, text) for pattern in patterns)


def _looks_like_section(text: str) -> bool:
    text = _clean_title(text)
    return bool(re.match(r"^(第[一二三四五六七八九十百\d]+节|\d+\.\d+)", text))
