"""教辅资料文本抽取（需求 T-K-02 / 7.2 第 2 步）。

按文件后缀选择解析器，抽取纯文本用于后续切分与向量化。
缺少对应解析库时优雅降级为空串，不影响主流程。
"""
from __future__ import annotations

import os


def detect_file_type(file_name: str) -> str:
    ext = os.path.splitext(file_name)[1].lower().lstrip(".")
    return {
        "pdf": "pdf",
        "doc": "word",
        "docx": "word",
        "ppt": "ppt",
        "pptx": "ppt",
        "txt": "txt",
        "md": "txt",
    }.get(ext, ext or "unknown")


def extract_text(file_path: str) -> str:
    """从本地文件抽取文本。返回可能为空。"""
    if not file_path or not os.path.exists(file_path):
        return ""
    ftype = detect_file_type(file_path)
    try:
        if ftype == "txt":
            return _read_txt(file_path)
        if ftype == "pdf":
            return _read_pdf(file_path)
        if ftype == "word":
            return _read_docx(file_path)
        if ftype == "ppt":
            return _read_pptx(file_path)
    except Exception:
        # 解析失败不阻断流程，交由上层标记失败或用占位
        return ""
    return ""


def _read_txt(path: str) -> str:
    for enc in ("utf-8", "gbk", "latin-1"):
        try:
            with open(path, encoding=enc) as f:
                return f.read()
        except UnicodeDecodeError:
            continue
    return ""


def _read_pdf(path: str) -> str:
    from pypdf import PdfReader

    reader = PdfReader(path)
    return "\n".join((page.extract_text() or "") for page in reader.pages)


def _read_docx(path: str) -> str:
    from docx import Document

    doc = Document(path)
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())


def _read_pptx(path: str) -> str:
    from pptx import Presentation

    prs = Presentation(path)
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                parts.append(shape.text_frame.text)
    return "\n".join(p for p in parts if p.strip())
